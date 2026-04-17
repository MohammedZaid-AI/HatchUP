import csv
import io
from typing import List, Dict, Union
import PyPDF2
from pptx import Presentation
from PIL import Image
import pytesseract
try:
    from docx import Document
except Exception:
    Document = None

class DocumentParser:
    """
    Handles extracting text from PDF, PPTX, and Image files.
    """
    
    @staticmethod
    def parse_file(uploaded_file) -> str:
        """
        Detects file type and delegates to the appropriate parser.
        Returns the extracted text as a single string.
        """
        filename = uploaded_file.name.lower()
        
        if filename.endswith(".pdf"):
            return DocumentParser._parse_pdf(uploaded_file)
        elif filename.endswith(".pptx") or filename.endswith(".ppt"):
            return DocumentParser._parse_pptx(uploaded_file)
        elif filename.endswith(".docx"):
            return DocumentParser._parse_docx(uploaded_file)
        elif filename.endswith(".txt"):
            return DocumentParser._parse_text(uploaded_file)
        elif filename.endswith(".csv"):
            return DocumentParser._parse_csv(uploaded_file)
        elif filename.endswith((".png", ".jpg", ".jpeg")):
            return DocumentParser._parse_image(uploaded_file)
        else:
            raise ValueError(f"Unsupported file format: {filename}")

    @staticmethod
    def _parse_pdf(file) -> str:
        text = ""
        try:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        except Exception as e:
            return f"Error parsing PDF: {str(e)}"
        return text

    @staticmethod
    def _parse_pptx(file) -> str:
        text = ""
        try:
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            return f"Error parsing PPTX: {str(e)}"
        return text

    @staticmethod
    def _parse_docx(file) -> str:
        if Document is None:
            return "Error parsing DOCX: python-docx is not installed."
        try:
            doc = Document(file)
            return "\n".join(paragraph.text for paragraph in doc.paragraphs if paragraph.text)
        except Exception as e:
            return f"Error parsing DOCX: {str(e)}"

    @staticmethod
    def _read_text_bytes(file) -> str:
        try:
            if hasattr(file, "seek"):
                file.seek(0)
            raw = file.read()
            if isinstance(raw, str):
                return raw
            return raw.decode("utf-8", errors="ignore")
        except Exception as e:
            return f"Error reading text file: {str(e)}"

    @staticmethod
    def _parse_text(file) -> str:
        return DocumentParser._read_text_bytes(file)

    @staticmethod
    def _parse_csv(file) -> str:
        try:
            text = DocumentParser._read_text_bytes(file)
            rows = csv.reader(io.StringIO(text))
            flattened = []
            for row in rows:
                cleaned = [cell.strip() for cell in row if cell and cell.strip()]
                if cleaned:
                    flattened.append(" | ".join(cleaned))
            return "\n".join(flattened)
        except Exception as e:
            return f"Error parsing CSV: {str(e)}"

    @staticmethod
    def _parse_image(file) -> str:
        """
        Uses Tesseract OCR to extract text from images.
        Requires Tesseract to be installed on the system.
        """
        try:
            image = Image.open(file)
            text = pytesseract.image_to_string(image)
            return text
        except Exception as e:
            return f"Error parsing Image (OCR): {str(e)}. Ensure Tesseract is installed."
